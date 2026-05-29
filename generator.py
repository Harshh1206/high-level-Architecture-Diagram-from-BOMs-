"""
generator.py — DeltaV Architecture PPT Drawing Engine  (v4 — Refactored)
=========================================================================
Generates ONE slide matching the DeltaV reference diagram style.

Uses ONLY the 9 hardware photos provided by the user:
  Charm.png · CIOC.png · DESKTOP.png · EWS-PROPlus.png
  PK_CONTROLLER.png · PRINTER.png · SWITCH.jpg · UPS.png · WORKSATION.png

All structural elements (rack frames, room borders, desk surface,
FOPP boxes, network bus) are drawn as vector shapes — no extra images.

IO Cabinet layout (top → bottom):
  ┌─────────────────────────┐
  │   PK_CONTROLLER.png     │  ← full-width, ~20% height
  │  ─── blue ctrl bus ───  │
  │  [CIOC][CIOC][CIOC]     │  ← CIOC.png portrait per tower (~18%)
  │  [CHRM][CHRM][CHRM]     │  ← Charm.png portrait per tower (~78%)
  └─────────────────────────┘
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Slide dimensions ──────────────────────────────────────────────────────────
SLIDE_W = 13.33
SLIDE_H = 7.50

# ── Image search directories ──────────────────────────────────────────────────
_BASE_DIR    = os.path.dirname(os.path.abspath(__file__))
IMAGES_DIR   = os.path.join(_BASE_DIR, 'images')
_SEARCH_DIRS = [IMAGES_DIR, os.path.join(_BASE_DIR, 'assets'), _BASE_DIR]

# ── ONLY the 9 user-provided hardware photos ──────────────────────────────────
_IMG_MAP: dict[str, list[str]] = {
    'CONTROLLER':   ['PK_CONTROLLER.png'],
    'CIOC':         ['CIOC.png'],
    'CHARM_BASE':   ['Charm.png'],
    'CHARM_AI':     ['Charm.png'],
    'CHARM_AO':     ['Charm.png'],
    'CHARM_DI':     ['Charm.png'],
    'CHARM_DO':     ['Charm.png'],
    'UPS':          ['UPS.png'],
    'WORKSTATION':  ['EWS-PROPlus.png'],
    'SWITCH':       ['SWITCH.jpg'],
    'MONITOR':      ['DESKTOP.png'],
    'OPERATOR_WS':  ['WORKSATION.png'],
    'PRINTER':      ['PRINTER.png'],
}

# Cabinet geometry constants
CAP_H             = 0.14
INNER_PAD         = 0.08
CHARMS_PER_TOWER  = 8
TOWERS_PER_IO_CAB = 3
CHARMS_PER_IO_CAB = 24


# ─────────────────────────────────────────────────────────────────────────────
# COLOUR PALETTE
# ─────────────────────────────────────────────────────────────────────────────

def rgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

C_CABINET   = rgb('252525')
C_CAB_INNER = rgb('E8E8E8')
C_RAIL_MED  = rgb('3C3C3C')
C_ROOM_BG   = rgb('F8F8F8')
C_OP_BG     = rgb('F8FBF8')
C_BORDER    = rgb('909090')
C_NET       = rgb('CC0077')
C_CTRL_BUS  = rgb('0055BB')
C_FOPP_FILL = rgb('CC0077')
C_TEXT      = rgb('1A1A1A')
C_TEXTMID   = rgb('555555')
C_WHITE     = rgb('FFFFFF')
C_SW_BG     = rgb('C8E6C9')
C_SW_BORDER = rgb('388E3C')
C_DESK_TOP  = rgb('7A6B50')
C_DESK_BODY = rgb('AAAAAA')


def _comp_color(dc: str) -> RGBColor:
    cmap = {
        'CONTROLLER': rgb('1A3A6B'), 'CIOC':       rgb('2E5FA3'),
        'CHARM_BASE': rgb('8B4513'), 'CHARM_AI':   rgb('C96A00'),
        'CHARM_AO':   rgb('B85C00'), 'CHARM_DI':   rgb('A85000'),
        'CHARM_DO':   rgb('963D00'), 'UPS':        rgb('B8860B'),
        'WORKSTATION':rgb('2B5EA7'), 'SWITCH':     rgb('2D6A2D'),
        'MONITOR':    rgb('4488AA'), 'OPERATOR_WS':rgb('1E4F9E'),
        'PRINTER':    rgb('888888'), 'SCREEN':     rgb('555566'),
        'FOPP':       rgb('7B3F9E'), 'FIREWALL':   rgb('C03030'),
        'MEDIA_CONV': rgb('00A0B0'), 'FIBER':      rgb('6040A0'),
        'POWER':      rgb('CC3300'), 'SOFTWARE':   rgb('4A8C3F'),
    }
    for k, v in cmap.items():
        if k in dc.upper():
            return v
    return rgb('888888')


# ─────────────────────────────────────────────────────────────────────────────
# LOW-LEVEL HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _img_path(key: str) -> str | None:
    for d in _SEARCH_DIRS:
        for f in _IMG_MAP.get(key, []):
            p = os.path.join(d, f)
            if os.path.exists(p):
                return p
    return None


def _rect(slide, x, y, w, h,
          fill=None, line_color=None, line_pt=0.75, no_fill=False):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if no_fill:
        s.fill.background()
    elif fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    s.text = ''
    return s


def _text(slide, text, x, y, w, h,
          size=9, bold=False, color=None, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text       = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.name  = 'Calibri'
    if color:
        r.font.color.rgb = color
    return tb


def _pic(slide, key, x, y, w, h, fallback=None):
    if w < 0.01 or h < 0.01:
        return None
    path = _img_path(key)
    if path:
        try:
            return slide.shapes.add_picture(
                path, Inches(x), Inches(y), Inches(w), Inches(h))
        except Exception:
            pass
    c = fallback or _comp_color(key)
    return _rect(slide, x, y, w, h, fill=c,
                 line_color=rgb('666666'), line_pt=0.4)


def _dash_border(shape, color, pt=1.2, dash='sysDash'):
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(pt)
    sp = shape._element.spPr
    le = sp.find(qn('a:ln'))
    if le is None:
        return
    etree.SubElement(le, qn('a:prstDash')).set('val', dash)


def _short(desc: str, qty: int) -> str:
    d = desc.strip()
    for a, b in [('DeltaV ', ''), ('Emerson ', ''), ('Assembly', 'Assy'),
                 ('Redundant', 'Red.'), ('Workstation', 'WS')]:
        d = d.replace(a, b)
    if len(d) > 34:
        d = d[:32] + '...'
    return (f'x{qty}  ' if qty > 1 else '') + d


# ─────────────────────────────────────────────────────────────────────────────
# RACK FRAME  — shapes only, no image
# ─────────────────────────────────────────────────────────────────────────────

def _draw_rack(slide, x, y, w, h):
    _rect(slide, x, y, w, h,
          fill=C_CABINET, line_color=rgb('111111'), line_pt=1.5)
    _rect(slide, x + INNER_PAD, y + CAP_H,
          w - 2 * INNER_PAD, h - 2 * CAP_H,
          fill=C_CAB_INNER, line_color=rgb('C0C0C0'), line_pt=0.4)
    _rect(slide, x + INNER_PAD, y + 0.030,
          w - 2 * INNER_PAD, 0.065, fill=C_RAIL_MED)
    _rect(slide, x + INNER_PAD, y + h - 0.095,
          w - 2 * INNER_PAD, 0.065, fill=C_RAIL_MED)
    _rect(slide, x + INNER_PAD, y + CAP_H,
          0.016, h - 2 * CAP_H, fill=rgb('404040'))
    _rect(slide, x + w - INNER_PAD - 0.016, y + CAP_H,
          0.016, h - 2 * CAP_H, fill=rgb('404040'))


# ─────────────────────────────────────────────────────────────────────────────
# SERVER CABINET
# ─────────────────────────────────────────────────────────────────────────────

_SERVER_H = {
    'SCREEN': 0.24, 'WORKSTATION': 0.50, 'UPS': 0.32,
    'SWITCH': 0.22, 'FIREWALL': 0.22, 'MEDIA_CONV': 0.20,
    'FIBER': 0.18, 'POWER': 0.20,
}


def _draw_server_cabinet(slide, x, y, w, h, items):
    _draw_rack(slide, x, y, w, h)
    cx    = x + INNER_PAD + 0.025
    cw    = w - 2 * INNER_PAD - 0.050
    cy    = y + CAP_H + 0.08
    avail = h - 2 * CAP_H - 0.16

    for item in items:
        dc = item.get('diagram_class', '')
        ih = next((v for k, v in _SERVER_H.items() if k in dc), 0.28)
        ih = min(ih, avail - (cy - (y + CAP_H + 0.08)) - 0.02)
        if ih < 0.09:
            break
        _pic(slide, dc, cx, cy, cw, ih - 0.025, _comp_color(dc))
        _text(slide, _short(item['description'], item['qty']),
              cx + 0.05, cy + 0.02, cw - 0.10, ih - 0.06,
              size=5.8, color=C_WHITE, wrap=False)
        cy += ih + 0.028


# ─────────────────────────────────────────────────────────────────────────────
# IO CABINET — redesigned layout matching reference diagram
# ─────────────────────────────────────────────────────────────────────────────

def _draw_io_cabinet(slide, x, y, w, h, items):
    _draw_rack(slide, x, y, w, h)

    cx      = x + INNER_PAD + 0.025
    cw      = w - 2 * INNER_PAD - 0.050
    inner_y = y + CAP_H + 0.07
    inner_h = h - 2 * CAP_H - 0.12

    ctrl_items  = [i for i in items if i['diagram_class'] in ('CONTROLLER', 'CIOC')]
    charm_items = [i for i in items if 'CHARM' in i.get('diagram_class', '')]

    cy           = inner_y
    ctrl_block_h = 0.0

    # ── PK Controller row ─────────────────────────────────────────────────────
    has_ctrl = any(i['diagram_class'] == 'CONTROLLER' for i in ctrl_items)
    if has_ctrl:
        pk_h = min(0.75, inner_h * 0.21)
        _pic(slide, 'CONTROLLER', cx, cy, cw, pk_h, _comp_color('CONTROLLER'))
        _text(slide, 'PK Controller',
              cx + 0.06, cy + pk_h * 0.28, cw - 0.12, 0.22,
              size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        cy          += pk_h + 0.04
        ctrl_block_h = pk_h + 0.04
        # Blue controller bus
        _rect(slide, cx + cw * 0.12, cy - 0.02, cw * 0.76, 0.018,
              fill=C_CTRL_BUS)
        cy += 0.025

    # ── CHARM tower columns ───────────────────────────────────────────────────
    bp_total = sum(
        i.get('qty', 1) for i in charm_items
        if 'CHARM' in i.get('diagram_class', '')
    )
    if bp_total == 0:
        bp_total = CHARMS_PER_TOWER

    n_towers     = min(max(1, math.ceil(bp_total / CHARMS_PER_TOWER)),
                       TOWERS_PER_IO_CAB)
    tower_area_h = max(0.5, inner_h - ctrl_block_h - 0.10)

    TOWER_GAP  = 0.030
    tower_w    = (cw - TOWER_GAP * (n_towers - 1)) / n_towers

    cioc_h  = tower_area_h * 0.18
    charm_h = tower_area_h * 0.78

    for t in range(n_towers):
        tx   = cx + t * (tower_w + TOWER_GAP)
        n_in = min(CHARMS_PER_TOWER, bp_total - t * CHARMS_PER_TOWER)
        if n_in <= 0:
            break

        # Blue drop from controller bus to CIOC
        if has_ctrl:
            drop_y = inner_y + ctrl_block_h - 0.022
            drop_h = cy - drop_y
            if drop_h > 0:
                _rect(slide, tx + tower_w / 2 - 0.008, drop_y,
                      0.016, drop_h, fill=C_CTRL_BUS)

        # CIOC image (portrait header)
        _pic(slide, 'CIOC', tx, cy, tower_w, cioc_h, _comp_color('CIOC'))
        _text(slide, f'CIOC {t + 1}',
              tx + 0.02, cy + cioc_h * 0.28, tower_w - 0.04, 0.18,
              size=5.8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # CHARM baseplate image (portrait, Charm.png)
        charm_y = cy + cioc_h + 0.012
        _pic(slide, 'CHARM_BASE', tx, charm_y, tower_w, charm_h,
             _comp_color('CHARM_BASE'))
        _text(slide, f'{n_in}/8',
              tx + 0.02, charm_y + charm_h * 0.46, tower_w - 0.04, 0.18,
              size=5.0, italic=True, color=rgb('F5DDB8'),
              align=PP_ALIGN.CENTER)

        # Tower separator
        if t < n_towers - 1:
            sx = tx + tower_w + TOWER_GAP / 2 - 0.004
            _rect(slide, sx, cy, 0.008, cioc_h + charm_h + 0.012,
                  fill=rgb('AAAAAA'))

    # Summary label
    label_y = cy + tower_area_h + 0.018
    _text(slide,
          f'{bp_total}/{n_towers * CHARMS_PER_TOWER} baseplates  ·  '
          f'{n_towers} CIOC tower{"s" if n_towers > 1 else ""}',
          cx, label_y, cw, 0.14,
          size=5.0, italic=True, color=C_TEXTMID, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SOFTWARE STRIP
# ─────────────────────────────────────────────────────────────────────────────

def _draw_software_strip(slide, x, y, w, h, items):
    _rect(slide, x, y, w, h, fill=C_SW_BG, line_color=C_SW_BORDER, line_pt=0.8)
    labels = []
    for it in items:
        d = it['description']
        for a, b in [('DeltaV ', ''), ('Emerson ', ''), ('Premium ', ''),
                     ('ProfessionalPLUS', 'ProPlus'), ('Software', 'SW')]:
            d = d.replace(a, b)
        labels.append(d[:26])
    _text(slide, '  |  '.join(labels),
          x + 0.10, y + 0.04, w - 0.20, h - 0.06,
          size=6.5, color=rgb('1B5E20'), wrap=False)


# ─────────────────────────────────────────────────────────────────────────────
# FOPP BOX  — magenta shape (matches reference)
# ─────────────────────────────────────────────────────────────────────────────

def _draw_fopp_box(slide, cx, cy, bw=0.36, bh=0.20):
    x = cx - bw / 2
    y = cy - bh / 2
    _rect(slide, x, y, bw, bh, fill=C_FOPP_FILL,
          line_color=rgb('880033'), line_pt=0.8)
    _text(slide, 'FOPP', x, y, bw, bh,
          size=6.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# OPERATOR DESK
# ─────────────────────────────────────────────────────────────────────────────

def _draw_operator_desk(slide, x, y, w, h, items):
    has_printer = any('PRINTER' in it['diagram_class'] for it in items)
    PR_W   = 0.60 if has_printer else 0.0
    DESK_W = w - PR_W - (0.08 if has_printer else 0.0)

    _text(slide, 'OWS',
          x + DESK_W * 0.20, y + 0.04, DESK_W * 0.60, 0.24,
          size=12, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)

    # Monitors — DESKTOP.png
    mon_h = min(h * 0.42, 1.65)
    mon_y = y + 0.32
    _pic(slide, 'MONITOR', x + 0.06, mon_y, DESK_W - 0.12, mon_h,
         _comp_color('MONITOR'))

    # Desk surface
    ds_y = mon_y + mon_h + 0.04
    ds_h = 0.10
    _rect(slide, x, ds_y, DESK_W, ds_h,
          fill=C_DESK_TOP, line_color=rgb('5A4A30'), line_pt=0.5)
    body_h = max(0.40, h - (ds_y + ds_h - y) - 0.10)
    _rect(slide, x, ds_y + ds_h, DESK_W, body_h,
          fill=C_DESK_BODY, line_color=rgb('888888'), line_pt=0.5)

    # Tower — WORKSATION.png
    twr_y = ds_y + ds_h + 0.07
    twr_h = max(0.40, body_h - 0.12)
    twr_w = min(DESK_W * 0.52, twr_h * 0.62)
    twr_x = x + (DESK_W - twr_w) / 2
    _pic(slide, 'OPERATOR_WS', twr_x, twr_y, twr_w, twr_h,
         _comp_color('OPERATOR_WS'))
    _text(slide, 'OWS Tower',
          twr_x, twr_y + twr_h + 0.02, twr_w, 0.14,
          size=6.0, color=C_TEXTMID, align=PP_ALIGN.CENTER)

    # Printer — PRINTER.png
    if has_printer:
        pr_x = x + DESK_W + 0.08
        pr_y = y + 0.32
        pr_h = min(0.55, h * 0.36)
        _pic(slide, 'PRINTER', pr_x, pr_y, PR_W - 0.06, pr_h,
             _comp_color('PRINTER'))
        _text(slide, 'Printer',
              pr_x, pr_y + pr_h + 0.03, PR_W - 0.06, 0.14,
              size=6.0, color=C_TEXTMID, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN SLIDE GENERATOR
# ─────────────────────────────────────────────────────────────────────────────

def generate_architecture_slide(
        structure: dict,
        project_title: str = 'DeltaV System Architecture') -> Presentation:

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb('EFEFEF')

    # ── Layout: centred diagram with generous margins ─────────────────────────
    # Horizontal: 1.50" margin each side  →  diagram width = 10.33"
    # Vertical  : 0.50" top margin (title), 1.42" bottom (legend box)
    H_OFFSET = 1.50          # left & right margin

    PDC_W    = 6.80           # PDC room width  (was 8.80)
    OP_GAP   = 0.12           # gap between rooms
    OP_W     = SLIDE_W - H_OFFSET - PDC_W - OP_GAP - H_OFFSET   # ≈ 3.41"

    PDC_X    = H_OFFSET       # = 1.50"
    OP_X     = PDC_X + PDC_W + OP_GAP

    MTOP     = 0.50           # space above rooms (title lives here)
    MBOT     = 1.42           # space below rooms (legend box lives here)
    ROOM_Y   = MTOP
    ROOM_H   = SLIDE_H - MTOP - MBOT   # ≈ 5.58"

    # ── Slide background ──────────────────────────────────────────────────────
    # Subtle off-white background behind the whole slide
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb('F2F2F2')

    # ── Title + underline (centred over diagram) ──────────────────────────────
    title_x = PDC_X
    title_w = PDC_W + OP_GAP + OP_W
    _text(slide, project_title,
          title_x, 0.08, title_w, 0.28,
          size=15, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
    _rect(slide, title_x + title_w * 0.20, 0.37,
          title_w * 0.60, 0.016, fill=C_TEXT)

    pdc_room = structure.get('PDC ROOM', {})
    op_room  = structure.get('OPERATOR ROOM', {})

    # ── Room fills + dashed borders ───────────────────────────────────────────
    _rect(slide, PDC_X, ROOM_Y, PDC_W, ROOM_H, fill=C_ROOM_BG)
    _rect(slide, OP_X,  ROOM_Y, OP_W,  ROOM_H, fill=C_OP_BG)

    pdc_b = _rect(slide, PDC_X, ROOM_Y, PDC_W, ROOM_H,
                  no_fill=True, line_color=C_BORDER, line_pt=1.2)
    _dash_border(pdc_b, C_BORDER, 1.2)
    _text(slide, 'PDC ROOM',
          PDC_X + 0.14, ROOM_Y + 0.07, 3.0, 0.22,
          size=11, bold=True, color=C_TEXT)

    op_b = _rect(slide, OP_X, ROOM_Y, OP_W, ROOM_H,
                 no_fill=True, line_color=C_BORDER, line_pt=1.2)
    _dash_border(op_b, C_BORDER, 1.2)
    _text(slide, 'OPERATOR ROOM',
          OP_X + 0.10, ROOM_Y + 0.07, OP_W - 0.18, 0.22,
          size=11, bold=True, color=C_TEXT)

    # ── Software strip ────────────────────────────────────────────────────────
    SW_H     = 0.23
    SW_Y     = ROOM_Y + 0.34
    sw_items = pdc_room.get('SOFTWARE', [])
    if sw_items:
        _draw_software_strip(slide, PDC_X + 0.13, SW_Y, PDC_W - 0.26, SW_H, sw_items)

    # ── Cabinet layout ────────────────────────────────────────────────────────
    pdc_cabs = {k: v for k, v in pdc_room.items()
                if k not in ('SOFTWARE', 'FOPP_NODES') and v}
    ordered  = sorted(pdc_cabs.keys(),
                      key=lambda k: (0 if 'SERVER' in k else 1, k))

    cab_top = (SW_Y + SW_H + 0.13) if sw_items else (ROOM_Y + 0.38)
    cab_h   = max(2.80, ROOM_H - (cab_top - ROOM_Y) - 0.80)
    n_cabs  = len(ordered)

    TOTAL_W = PDC_W - 0.28
    GAP     = 0.12
    if n_cabs > 0:
        s_frac  = 0.22
        io_frac = (1 - s_frac) / max(n_cabs - 1, 1)
        widths  = [s_frac if 'SERVER' in k else io_frac for k in ordered]
        tot     = sum(widths)
        widths  = [ww / tot for ww in widths]
    else:
        widths = []

    usable = TOTAL_W - GAP * max(0, n_cabs - 1)
    cab_x  = PDC_X + 0.14
    mid_xs = []

    for i, cname in enumerate(ordered):
        items = pdc_cabs[cname]
        cw    = usable * widths[i]
        mid_xs.append(cab_x + cw / 2)

        if 'SERVER' in cname:
            _draw_server_cabinet(slide, cab_x, cab_top, cw, cab_h, items)
        else:
            _draw_io_cabinet(slide, cab_x, cab_top, cw, cab_h, items)

        _text(slide, cname.replace('_', ' ').title(),
              cab_x, cab_top + cab_h + 0.04, cw, 0.17,
              size=7.5, color=C_TEXTMID, align=PP_ALIGN.CENTER)
        cab_x += cw + GAP

    # ── Network bus (magenta horizontal + vertical drops) ─────────────────────
    NET_Y  = cab_top - 0.12
    net_x1 = PDC_X + 0.14 + 0.04
    net_x2 = cab_x - GAP - 0.04
    if mid_xs:
        _rect(slide, net_x1, NET_Y, net_x2 - net_x1, 0.022, fill=C_NET)
        for mx in mid_xs:
            _rect(slide, mx - 0.009, NET_Y, 0.018, cab_top - NET_Y, fill=C_NET)

    # ── FOPP boxes below cabinets ─────────────────────────────────────────────
    FOPP_Y     = min(cab_top + cab_h + 0.38,
                     ROOM_Y + ROOM_H - 0.36)   # never overflows room
    fopp_items = pdc_room.get('FOPP_NODES', [])
    n_fopp     = max(1, len(fopp_items))
    spacing    = (PDC_W * 0.55) / (n_fopp + 1)
    pdc_fxs    = [PDC_X + PDC_W * 0.45 + spacing * (i + 1) for i in range(n_fopp)]

    for i, fx in enumerate(pdc_fxs):
        _rect(slide, fx - 0.009, NET_Y, 0.018, FOPP_Y - NET_Y + 0.07,
              fill=C_NET)
        _draw_fopp_box(slide, fx, FOPP_Y)
        lbl = fopp_items[i]['description'] if i < len(fopp_items) else 'FOPP'
        _text(slide, lbl[:14],
              fx - 0.28, FOPP_Y - 0.26, 0.56, 0.17,
              size=5.8, bold=True, color=C_TEXTMID, align=PP_ALIGN.CENTER)

    pdc_fx = pdc_fxs[-1]

    # ── Operator FOPP ─────────────────────────────────────────────────────────
    op_fx = OP_X + 0.42
    _rect(slide, op_fx - 0.009, ROOM_Y + 0.52, 0.018,
          FOPP_Y - (ROOM_Y + 0.52), fill=C_NET)
    _draw_fopp_box(slide, op_fx, FOPP_Y)
    _text(slide, 'FOPP',
          op_fx - 0.20, FOPP_Y - 0.26, 0.40, 0.17,
          size=5.8, bold=True, color=C_TEXTMID, align=PP_ALIGN.CENTER)

    # ── Cross-room fibre link ─────────────────────────────────────────────────
    lx1 = pdc_fx + 0.18
    lx2 = op_fx  - 0.18
    if lx2 > lx1:
        _rect(slide, lx1, FOPP_Y - 0.009, lx2 - lx1, 0.018, fill=C_NET)

    # ── Operator desk ─────────────────────────────────────────────────────────
    desk_items = op_room.get('OPERATOR_DESK', [])
    dx = OP_X + 0.09
    dy = ROOM_Y + 0.34
    dw = OP_W - 0.18
    dh = max(3.10, min(cab_top + cab_h - dy + 0.20, ROOM_H - 0.50))
    _draw_operator_desk(slide, dx, dy, dw, dh, desk_items)

    # ── Legend box — bottom right ─────────────────────────────────────────────
    legend_entries = [
        ('PK Controller', rgb('1A3A6B')),
        ('CIOC',          rgb('2E5FA3')),
        ('CHARM I/O',     rgb('CC6600')),
        ('EWS / WS',      rgb('2B5EA7')),
        ('Switch',        rgb('2D6A2D')),
        ('UPS',           rgb('B8860B')),
        ('Fiber / FOPP',  rgb('7B3F9E')),
        ('Network',       C_NET),
        ('Software',      rgb('4A8C3F')),
        ('Power',         rgb('CC3300')),
    ]

    COLS      = 2
    PER_COL   = math.ceil(len(legend_entries) / COLS)
    ROW_H     = 0.195
    SWATCH_W  = 0.115
    SWATCH_H  = 0.130
    LB_PAD_X  = 0.13
    LB_PAD_Y  = 0.25   # top padding inside box (below header)
    LB_W      = 2.80
    LB_H      = LB_PAD_Y + PER_COL * ROW_H + 0.12

    # Align box flush with right edge of diagram, sitting in the bottom margin
    LB_X = OP_X + OP_W - LB_W   # right edge of diagram
    LB_Y = ROOM_Y + ROOM_H + 0.16

    # Box background + border
    _rect(slide, LB_X, LB_Y, LB_W, LB_H,
          fill=rgb('FFFFFF'), line_color=rgb('888888'), line_pt=0.9)

    # "LEGEND" header bar
    _rect(slide, LB_X, LB_Y, LB_W, 0.20,
          fill=rgb('2C2C2C'), line_color=None)
    _text(slide, 'LEGEND',
          LB_X + 0.10, LB_Y + 0.02, LB_W - 0.20, 0.17,
          size=7, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    COL_W = (LB_W - LB_PAD_X * 2) / COLS

    for idx, (lbl, col) in enumerate(legend_entries):
        ci    = idx // PER_COL          # column index
        ri    = idx %  PER_COL          # row index within column
        ex    = LB_X + LB_PAD_X + ci * COL_W
        ey    = LB_Y + LB_PAD_Y + ri * ROW_H

        # Colour swatch
        _rect(slide, ex, ey, SWATCH_W, SWATCH_H,
              fill=col, line_color=rgb('AAAAAA'), line_pt=0.3)

        # Label
        _text(slide, lbl,
              ex + SWATCH_W + 0.07, ey - 0.01,
              COL_W - SWATCH_W - 0.09, ROW_H,
              size=6.0, color=C_TEXT)

    return prs


def generate_pptx(structure: dict, output_path: str,
                  project_title: str = 'DeltaV System Architecture') -> str:
    prs = generate_architecture_slide(structure, project_title)
    prs.save(output_path)
    return output_path
"""
app.py — DeltaV Architecture Generator  (v3 — Image Edition)
=============================================================
Flow:
  1. Upload BOM  →  parse + classify automatically
  2. Review table; correct any UNKNOWNs via dropdowns
  3. Architecture diagram auto-generated and shown as PNG preview
  4. One-click PPTX download

No "Generate" button required — the diagram is always in sync with the BOM.
"""

import os, sys, io, tempfile, subprocess, shutil
import pandas as pd
import streamlit as st

sys.path.insert(0, os.path.dirname(__file__))
from parser     import parse_bom
from classifier import classify_dataframe, load_rules, save_user_correction
from grouper    import group_bom
from generator  import generate_pptx

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title='DeltaV Architecture Generator',
    page_icon='⚡',
    layout='wide',
    initial_sidebar_state='expanded',
)

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=JetBrains+Mono:wght@300;400;500&family=Inter:wght@300;400;500&display=swap');
*,*::before,*::after{box-sizing:border-box;}
html,body,[data-testid="stAppViewContainer"]{background:#0B0E14 !important;color:#E8EAF0 !important;font-family:'Inter',sans-serif !important;}
[data-testid="stAppViewContainer"]{background:radial-gradient(ellipse 80% 60% at 50% -10%,rgba(255,160,30,.07) 0%,transparent 70%),repeating-linear-gradient(0deg,transparent,transparent 39px,rgba(255,255,255,.018) 39px,rgba(255,255,255,.018) 40px),repeating-linear-gradient(90deg,transparent,transparent 39px,rgba(255,255,255,.018) 39px,rgba(255,255,255,.018) 40px),#0B0E14 !important;background-attachment:fixed !important;}
[data-testid="stSidebar"]{background:#0D1018 !important;border-right:1px solid rgba(255,160,30,.15) !important;}
[data-testid="stSidebar"] *{color:#B8BCC8 !important;}
#MainMenu,footer,header{visibility:hidden;}
.block-container{padding:0 2rem 2rem 2rem !important;max-width:1400px !important;}
h1,h2,h3{font-family:'Syne',sans-serif !important;}
.hero{padding:3rem 0 2.5rem 0;}
.hero-eyebrow{font-family:'JetBrains Mono',monospace;font-size:11px;font-weight:500;letter-spacing:.22em;color:#FFA01E;text-transform:uppercase;margin-bottom:.8rem;}
.hero-title{font-family:'Syne',sans-serif;font-size:clamp(2rem,5vw,3.6rem);font-weight:800;line-height:1.05;color:#F0F2F8;margin:0 0 1rem 0;letter-spacing:-.02em;}
.hero-title .accent{color:#FFA01E;}
.hero-sub{font-size:15px;color:#7A8099;font-weight:300;max-width:580px;line-height:1.7;}
.hero-divider{height:1px;background:linear-gradient(90deg,#FFA01E 0%,rgba(255,160,30,.2) 40%,transparent 100%);margin:2.5rem 0 2rem 0;}
.metrics-row{display:grid;grid-template-columns:repeat(4,1fr);gap:1rem;margin-bottom:2rem;}
.metric-card{background:rgba(255,255,255,.033);border:1px solid rgba(255,255,255,.07);border-radius:8px;padding:1.2rem 1.4rem;position:relative;overflow:hidden;}
.metric-card::before{content:'';position:absolute;top:0;left:0;right:0;height:2px;background:linear-gradient(90deg,#FFA01E,transparent);}
.metric-label{font-family:'JetBrains Mono',monospace;font-size:10px;letter-spacing:.15em;color:#5A6080;text-transform:uppercase;margin-bottom:.5rem;}
.metric-value{font-family:'Syne',sans-serif;font-size:2.2rem;font-weight:700;color:#F0F2F8;line-height:1;}
.metric-value.amber{color:#FFA01E;} .metric-value.green{color:#4ADE80;} .metric-value.red{color:#F87171;}
.metric-sub{font-size:11px;color:#5A6080;margin-top:.3rem;font-family:'JetBrains Mono',monospace;}
.section-header{display:flex;align-items:center;gap:1rem;margin-bottom:1rem;}
.section-num{font-family:'JetBrains Mono',monospace;font-size:10px;color:#FFA01E;letter-spacing:.1em;background:rgba(255,160,30,.08);border:1px solid rgba(255,160,30,.2);border-radius:4px;padding:3px 8px;flex-shrink:0;}
.section-title{font-family:'Syne',sans-serif;font-size:18px;font-weight:700;color:#E8EAF0;}
.section-line{flex:1;height:1px;background:rgba(255,255,255,.06);}
[data-testid="stFileUploader"]>div{border:1.5px dashed rgba(255,160,30,.25) !important;background:rgba(255,160,30,.025) !important;border-radius:12px !important;padding:1.5rem !important;}
[data-testid="stFileUploader"]>div:hover{border-color:rgba(255,160,30,.55) !important;background:rgba(255,160,30,.05) !important;}
[data-testid="stDataFrame"] th{background:rgba(255,160,30,.08) !important;color:#FFA01E !important;font-family:'JetBrains Mono',monospace !important;font-size:10px !important;letter-spacing:.1em !important;text-transform:uppercase !important;border-bottom:1px solid rgba(255,160,30,.2) !important;padding:10px 14px !important;}
[data-testid="stDataFrame"] td{background:rgba(255,255,255,.02) !important;color:#B8BCC8 !important;border-bottom:1px solid rgba(255,255,255,.04) !important;font-family:'JetBrains Mono',monospace !important;font-size:11px !important;padding:8px 14px !important;}
[data-testid="stButton"] button{background:linear-gradient(135deg,#FFA01E 0%,#FF7A00 100%) !important;color:#0B0E14 !important;border:none !important;border-radius:8px !important;font-family:'Syne',sans-serif !important;font-size:14px !important;font-weight:700 !important;letter-spacing:.04em !important;box-shadow:0 4px 24px rgba(255,160,30,.25) !important;text-transform:uppercase !important;}
[data-testid="stButton"] button:hover{transform:translateY(-1px) !important;box-shadow:0 8px 32px rgba(255,160,30,.4) !important;}
[data-testid="stDownloadButton"] button{background:rgba(74,222,128,.1) !important;color:#4ADE80 !important;border:1px solid rgba(74,222,128,.3) !important;border-radius:8px !important;font-family:'Syne',sans-serif !important;font-size:14px !important;font-weight:700 !important;width:100% !important;text-transform:uppercase !important;}
[data-testid="stDownloadButton"] button:hover{background:rgba(74,222,128,.18) !important;box-shadow:0 4px 20px rgba(74,222,128,.2) !important;}
[data-testid="stSelectbox"]>div{background:rgba(255,255,255,.04) !important;border:1px solid rgba(255,255,255,.1) !important;border-radius:6px !important;color:#E8EAF0 !important;}
[data-testid="stTextInput"] input{background:rgba(255,255,255,.04) !important;border:1px solid rgba(255,255,255,.1) !important;border-radius:6px !important;color:#E8EAF0 !important;font-family:'JetBrains Mono',monospace !important;font-size:13px !important;}
.preview-wrap{position:relative;border-radius:12px;overflow:hidden;border:1px solid rgba(255,160,30,.18);box-shadow:0 0 0 1px rgba(255,255,255,.04),0 24px 64px rgba(0,0,0,.6),0 0 80px rgba(255,160,30,.06);background:#141720;margin-bottom:1.5rem;}
.preview-topbar{display:flex;align-items:center;justify-content:space-between;padding:.65rem 1rem;background:rgba(255,255,255,.03);border-bottom:1px solid rgba(255,255,255,.06);}
.preview-dots{display:flex;gap:6px;}
.preview-dot{width:10px;height:10px;border-radius:50%;}
.preview-dot-r{background:#FF5F57;} .preview-dot-y{background:#FFBD2E;} .preview-dot-g{background:#28CA41;}
.preview-label{font-family:'JetBrains Mono',monospace;font-size:10px;color:#404560;letter-spacing:.1em;}
.preview-badge{font-family:'JetBrains Mono',monospace;font-size:9px;color:#4ADE80;background:rgba(74,222,128,.08);border:1px solid rgba(74,222,128,.2);border-radius:4px;padding:2px 7px;letter-spacing:.08em;}
.preview-footer{padding:.6rem 1rem;background:rgba(255,255,255,.02);border-top:1px solid rgba(255,255,255,.04);display:flex;justify-content:space-between;align-items:center;}
.preview-footer-left{font-family:'JetBrains Mono',monospace;font-size:10px;color:#2A2E40;}
.sidebar-brand{font-family:'Syne',sans-serif;font-size:20px;font-weight:800;letter-spacing:-.02em;}
.sidebar-ver{font-family:'JetBrains Mono',monospace;font-size:10px;letter-spacing:.1em;color:#404560;margin-bottom:1.5rem;}
.sidebar-sec{font-family:'JetBrains Mono',monospace;font-size:9px;letter-spacing:.18em;text-transform:uppercase;color:#404560;margin:1.2rem 0 .6rem 0;padding-bottom:.4rem;border-bottom:1px solid rgba(255,255,255,.05);}
.tag-pill{display:inline-block;font-family:'JetBrains Mono',monospace;font-size:10px;padding:2px 8px;border-radius:20px;border:1px solid rgba(255,160,30,.3);color:#FFA01E;background:rgba(255,160,30,.07);margin:2px;}
::-webkit-scrollbar{width:6px;height:6px;}
::-webkit-scrollbar-track{background:transparent;}
::-webkit-scrollbar-thumb{background:rgba(255,160,30,.2);border-radius:3px;}
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# PPTX → PNG  (LibreOffice → Pillow fallback)
# ══════════════════════════════════════════════════════════════════════════════

def pptx_to_preview_image(pptx_bytes: bytes) -> bytes | None:
    # PATH 1 — LibreOffice + pdftoppm
    try:
        import socket as _socket
        from pathlib import Path as _Path
        if not shutil.which('soffice'):
            raise RuntimeError('no soffice')
        _tmp = tempfile.mkdtemp()
        _pptx = os.path.join(_tmp, 'slide.pptx')
        _pdf  = os.path.join(_tmp, 'slide.pdf')
        with open(_pptx, 'wb') as f:
            f.write(pptx_bytes)
        _env = os.environ.copy()
        _env.update({'SAL_USE_VCLPLUGIN': 'svp', 'HOME': _tmp})

        def _needs_shim():
            try:
                s = _socket.socket(_socket.AF_UNIX, _socket.SOCK_STREAM); s.close(); return False
            except OSError:
                return True

        if _needs_shim():
            _shim = _Path(tempfile.gettempdir()) / 'lo_shim.so'
            if not _shim.exists():
                _src = _Path(tempfile.gettempdir()) / 'lo_shim.c'
                _src.write_text('#define _GNU_SOURCE\n#include <dlfcn.h>\n#include <errno.h>\n#include <sys/socket.h>\nstatic int (*rs)(int,int,int);static int (*rsp)(int,int,int,int[2]);\n__attribute__((constructor)) static void i(void){rs=dlsym(RTLD_NEXT,"socket");rsp=dlsym(RTLD_NEXT,"socketpair");}\nint socket(int d,int t,int p){if(d==AF_UNIX){int fd=rs(d,t,p);if(fd>=0)return fd;int sv[2];if(rsp(d,t,p,sv)==0)return sv[0];errno=EPERM;return -1;}return rs(d,t,p);}')
                subprocess.run(['gcc','-shared','-fPIC','-o',str(_shim),str(_src),'-ldl'],check=True,capture_output=True)
                _src.unlink()
            _env['LD_PRELOAD'] = str(_shim)

        r = subprocess.run(['soffice','--headless','--convert-to','pdf','--outdir',_tmp,_pptx],
                           env=_env, capture_output=True, timeout=60)
        if r.returncode != 0 or not os.path.exists(_pdf):
            raise RuntimeError('soffice failed')
        if not shutil.which('pdftoppm'):
            raise RuntimeError('no pdftoppm')
        subprocess.run(['pdftoppm','-jpeg','-r','180','-f','1','-l','1',_pdf,os.path.join(_tmp,'pg')],
                       check=True, capture_output=True, timeout=30)
        imgs = sorted(_Path(_tmp).glob('pg-*.jpg'))
        if not imgs:
            raise RuntimeError('no img')
        raw = imgs[0].read_bytes()
        shutil.rmtree(_tmp, ignore_errors=True)
        from PIL import Image as _PI
        pil = _PI.open(io.BytesIO(raw)).convert('RGB')
        buf = io.BytesIO(); pil.save(buf,'PNG',optimize=True); return buf.getvalue()
    except Exception:
        pass

    # PATH 2 — Pillow shape renderer with image support
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image, ImageDraw, ImageFont
        import textwrap as _tw
        SC = 150
        CW, CH = int(13.33*SC), int(7.50*SC)
        prs   = Presentation(io.BytesIO(pptx_bytes))
        slide = prs.slides[0]
        img   = Image.new('RGB',(CW,CH),(255,255,255))
        draw  = ImageDraw.Draw(img)
        def ep(e): return int((e or 0)/914400*SC)
        _fc: dict = {}
        _FP = ['/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
               '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf']
        def _f(sz):
            if sz not in _fc:
                for p in _FP:
                    try: _fc[sz]=ImageFont.truetype(p,max(6,sz)); break
                    except: pass
                else: _fc[sz]=ImageFont.load_default()
            return _fc[sz]
        try:
            bg = slide.background.fill.fore_color.rgb
            img.paste(Image.new('RGB',(CW,CH),(bg[0],bg[1],bg[2])))
        except: pass
        for sh in slide.shapes:
            try:
                x,y,w,h = ep(sh.left),ep(sh.top),ep(sh.width),ep(sh.height)
                if w<=0 or h<=0: continue
                
                # ── Handle picture shapes (images) ──
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        pic_img = Image.open(io.BytesIO(sh.image.blob)).convert('RGB')
                        pic_resized = pic_img.resize((max(1, int(w)), max(1, int(h))), Image.Resampling.LANCZOS)
                        img.paste(pic_resized, (max(0, int(x)), max(0, int(y))))
                    except Exception:
                        pass
                
                # ── Handle filled shapes ──
                try:
                    fi = sh.fill
                    if fi.type is not None:
                        c = fi.fore_color.rgb; draw.rectangle([x,y,x+w,y+h],fill=(c[0],c[1],c[2]))
                except: pass
                
                # ── Handle text ──
                if sh.has_text_frame:
                    txt = sh.text_frame.text.strip()
                    if txt:
                        try: cr=sh.text_frame.paragraphs[0].runs[0].font.color.rgb; tc=(cr[0],cr[1],cr[2])
                        except: tc=(220,220,220)
                        fo = _f(9)
                        for i,line in enumerate(_tw.wrap(txt,width=max(3,w//6)) or [txt]):
                            try:
                                bb=draw.textbbox((0,0),line,font=fo)
                                draw.text((x+max(0,(w-(bb[2]-bb[0]))//2),y+4+i*11),line,fill=tc,font=fo)
                            except: pass
            except: continue
        buf=io.BytesIO(); img.save(buf,'PNG',optimize=True); return buf.getvalue()
    except:
        return None


# ══════════════════════════════════════════════════════════════════════════════
# CORE PIPELINE  (parse → classify → group → pptx → preview)
# ══════════════════════════════════════════════════════════════════════════════

def _generate(df: pd.DataFrame, corrections: dict, project_title: str) -> None:
    """Apply corrections, run grouper, generate PPTX + PNG, cache in session."""
    for idx, cls in corrections.items():
        df.at[idx, 'diagram_class'] = cls
    df.loc[df['diagram_class'] == 'UNKNOWN', 'diagram_class'] = 'WORKSTATION'

    structure  = group_bom(df)
    out        = tempfile.mktemp(suffix='.pptx')
    generate_pptx(structure, out, project_title=project_title)
    with open(out, 'rb') as f:
        pptx_bytes = f.read()
    os.unlink(out)

    st.session_state['pptx_bytes']    = pptx_bytes
    st.session_state['preview_img']   = pptx_to_preview_image(pptx_bytes)
    st.session_state['preview_title'] = project_title


# ══════════════════════════════════════════════════════════════════════════════
# SIDEBAR
# ══════════════════════════════════════════════════════════════════════════════

with st.sidebar:
    st.markdown('<div class="sidebar-brand" style="color:#FFA01E">⚡ DeltaV</div>', unsafe_allow_html=True)
    st.markdown('<div style="color:#C0C4D8;font-family:\'Syne\',sans-serif;font-size:13px;font-weight:600;margin-bottom:2px;">Architecture Generator</div>', unsafe_allow_html=True)
    st.markdown('<div class="sidebar-ver">v3.0 · IMAGE EDITION</div>', unsafe_allow_html=True)

    st.markdown('<div class="sidebar-sec">Project Settings</div>', unsafe_allow_html=True)
    project_title = st.text_input('PROJECT TITLE', value='CO2 Capture Plant')

    st.markdown('<div class="sidebar-sec">Supported Formats</div>', unsafe_allow_html=True)
    st.markdown('<span class="tag-pill">.xlsx</span><span class="tag-pill">.xls</span><span class="tag-pill">.csv</span>', unsafe_allow_html=True)

    st.markdown('<div class="sidebar-sec">Image Rendering</div>', unsafe_allow_html=True)
    try:
        from generator import IMAGES_DIR
        n = len([f for f in os.listdir(IMAGES_DIR) if f.lower().endswith(('.png','.jpg'))]) if os.path.isdir(IMAGES_DIR) else 0
        st.markdown(f'<div style="font-family:\'JetBrains Mono\',monospace;font-size:11px;color:#5A6080;line-height:1.8;">{n} component images<br>→ <code>./images/</code></div>', unsafe_allow_html=True)
    except Exception:
        pass

    st.markdown('<div class="sidebar-sec">Engine</div>', unsafe_allow_html=True)
    st.markdown('<div style="font-size:12px;color:#5A6080;line-height:1.7;">Rule-based · 100% offline<br>Self-learning corrections</div>', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# HERO
# ══════════════════════════════════════════════════════════════════════════════

st.markdown("""
<div class="hero">
    <div class="hero-eyebrow">⬡  Automation Solutions · Pune · DCS Engineering</div>
    <h1 class="hero-title">BOM <span class="accent">→</span> Architecture<br>in Seconds</h1>
    <p class="hero-sub">Upload any Bill of Materials — the rule engine classifies every component, groups it into cabinets, and instantly renders a DeltaV architecture diagram using real hardware images.</p>
</div>
<div class="hero-divider"></div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# STEP 1 — UPLOAD
# ══════════════════════════════════════════════════════════════════════════════

st.markdown('<div class="section-header"><span class="section-num">01</span><span class="section-title">Upload Bill of Materials</span><div class="section-line"></div></div>', unsafe_allow_html=True)

uploaded_file = st.file_uploader(
    'Drop your BOM file — Excel or CSV, any column order',
    type=['xlsx','xls','csv'],
)

if uploaded_file is None:
    st.markdown("""
    <div style="margin-top:1rem;padding:1.5rem;background:rgba(255,255,255,.02);border:1px solid rgba(255,255,255,.05);border-radius:8px;">
        <div style="font-family:'JetBrains Mono',monospace;font-size:11px;color:#404560;line-height:2.2;">
        ◦ &nbsp;Any column order — auto-detects Description, Qty, Area, Part No<br>
        ◦ &nbsp;Headers on any row (1–15), merged cells, section labels handled<br>
        ◦ &nbsp;UTF-8, Latin-1, Windows-1252 encodings supported<br>
        ◦ &nbsp;Multi-sheet Excel — best sheet auto-selected
        </div>
    </div>
    """, unsafe_allow_html=True)
    st.stop()


# ══════════════════════════════════════════════════════════════════════════════
# PARSE + CLASSIFY
# ══════════════════════════════════════════════════════════════════════════════

suffix = os.path.splitext(uploaded_file.name)[1]
with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
    tmp.write(uploaded_file.read())
    tmp_path = tmp.name

with st.spinner('Parsing and classifying…'):
    try:
        df = parse_bom(tmp_path)
    except Exception as e:
        st.error(f'**Parse error:** {e}')
        os.unlink(tmp_path)
        st.stop()
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

    rules = load_rules()
    df    = classify_dataframe(df, rules)

unknowns = df[df['diagram_class'] == 'UNKNOWN']
known    = df[df['diagram_class'] != 'UNKNOWN']
pct      = int(len(known)/len(df)*100) if len(df) else 0

# Metrics
st.markdown(f"""
<div class="metrics-row">
    <div class="metric-card">
        <div class="metric-label">Total Items</div>
        <div class="metric-value">{len(df)}</div>
        <div class="metric-sub">{uploaded_file.name[:22]}</div>
    </div>
    <div class="metric-card">
        <div class="metric-label">Classified</div>
        <div class="metric-value green">{len(known)}</div>
        <div class="metric-sub">{pct}% auto-resolved</div>
    </div>
    <div class="metric-card">
        <div class="metric-label">Unknowns</div>
        <div class="metric-value {'red' if len(unknowns) else 'green'}">{len(unknowns)}</div>
        <div class="metric-sub">{'need review' if len(unknowns) else 'all clear ✓'}</div>
    </div>
    <div class="metric-card">
        <div class="metric-label">PDC Items</div>
        <div class="metric-value amber">{len(df[df['area']=='PDC ROOM'])}</div>
        <div class="metric-sub">{len(df[df['area']=='OPERATOR ROOM'])} operator room</div>
    </div>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# STEP 2 — REVIEW CLASSIFICATION
# ══════════════════════════════════════════════════════════════════════════════

st.markdown('<div class="section-header"><span class="section-num">02</span><span class="section-title">Review Classification</span><div class="section-line"></div></div>', unsafe_allow_html=True)

disp = df[['sr_no','area','description','qty','diagram_class','confidence','part_number']].copy()
disp.columns = ['SR','AREA','DESCRIPTION','QTY','CLASS','CONFIDENCE','PART NO']
st.dataframe(disp, use_container_width=True, height=280)

# Manual corrections
corrections: dict[int, str] = {}
if len(unknowns):
    st.markdown(f'<div style="font-family:\'Syne\',sans-serif;font-size:15px;font-weight:600;color:#F87171;margin:1.2rem 0 .8rem;">⚠ {len(unknowns)} item{"s" if len(unknowns)>1 else ""} need manual classification</div>', unsafe_allow_html=True)
    all_classes = sorted(rules.keys())
    for idx, row in unknowns.iterrows():
        c1, c2 = st.columns([3,1])
        with c1:
            st.markdown(f'<div style="padding:.5rem 0;font-size:13px;color:#E8EAF0;">{row["description"][:80]}<br><span style="font-family:JetBrains Mono,monospace;font-size:10px;color:#404560;">{row.get("part_number","")}</span></div>', unsafe_allow_html=True)
        with c2:
            choice = st.selectbox('',['— select —']+all_classes,key=f'fix_{idx}',label_visibility='collapsed')
            if choice != '— select —':
                corrections[idx] = choice

    if corrections and st.button('💾  Save corrections to rules.json'):
        for idx, cls in corrections.items():
            save_user_correction(df.loc[idx,'description'], cls)
        st.success(f'✅ Saved {len(corrections)} corrections — auto-classifies next time')
else:
    st.success('✅  All items classified automatically — no review needed')


# ══════════════════════════════════════════════════════════════════════════════
# AUTO-GENERATE  (re-runs whenever BOM or corrections change)
# ══════════════════════════════════════════════════════════════════════════════

_sig = (uploaded_file.name, len(df), tuple(sorted(corrections.items())), project_title)
if st.session_state.get('_gen_sig') != _sig:
    with st.spinner('Generating architecture diagram…'):
        _generate(df.copy(), corrections, project_title)
    st.session_state['_gen_sig'] = _sig


# ══════════════════════════════════════════════════════════════════════════════
# STEP 3 — PREVIEW + DOWNLOAD
# ══════════════════════════════════════════════════════════════════════════════

st.markdown('<div class="section-header" style="margin-top:2rem;"><span class="section-num">03</span><span class="section-title">Architecture Diagram</span><div class="section-line"></div></div>', unsafe_allow_html=True)

pptx_bytes = st.session_state.get('pptx_bytes')
img_bytes  = st.session_state.get('preview_img')

if pptx_bytes:
    fname = st.session_state.get('preview_title', project_title).replace(' ','_').replace('/','_') + '_Architecture.pptx'

    # Preview frame
    st.markdown("""
    <div class="preview-wrap">
        <div class="preview-topbar">
            <div class="preview-dots">
                <div class="preview-dot preview-dot-r"></div>
                <div class="preview-dot preview-dot-y"></div>
                <div class="preview-dot preview-dot-g"></div>
            </div>
            <div class="preview-label">ARCHITECTURE DIAGRAM · SLIDE 1 OF 1</div>
            <div class="preview-badge">✓ GENERATED</div>
        </div>
        <div style="background:#0D0F16;padding:0;">
    """, unsafe_allow_html=True)

    if img_bytes:
        st.image(img_bytes, use_column_width=True)
    else:
        st.markdown('<div style="padding:3rem;text-align:center;font-family:\'JetBrains Mono\',monospace;font-size:12px;color:#404560;">Preview unavailable — download the PPTX to view</div>', unsafe_allow_html=True)

    st.markdown(f"""
        </div>
        <div class="preview-footer">
            <div class="preview-footer-left">DELTAV SYSTEM ARCHITECTURE · {st.session_state.get('preview_title','').upper()}</div>
            <div style="font-family:'JetBrains Mono',monospace;font-size:10px;color:#2A2E40;">DOWNLOAD FOR FULL RESOLUTION</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # Download + file info
    col_dl, col_info = st.columns([2,3])
    with col_dl:
        st.download_button(
            '↓  DOWNLOAD ARCHITECTURE.PPTX',
            data=pptx_bytes,
            file_name=fname,
            mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            use_container_width=True,
        )
    with col_info:
        st.markdown(f"""
        <div style="padding:.7rem 1rem;background:rgba(255,255,255,.02);border:1px solid rgba(255,255,255,.06);border-radius:8px;display:flex;gap:2rem;align-items:center;">
            <div><div style="font-family:'JetBrains Mono',monospace;font-size:9px;letter-spacing:.12em;color:#404560;text-transform:uppercase;margin-bottom:3px;">Format</div>
                 <div style="font-family:'Syne',sans-serif;font-size:13px;font-weight:600;color:#C0C4D8;">.pptx · Microsoft PowerPoint</div></div>
            <div><div style="font-family:'JetBrains Mono',monospace;font-size:9px;letter-spacing:.12em;color:#404560;text-transform:uppercase;margin-bottom:3px;">Items</div>
                 <div style="font-family:'Syne',sans-serif;font-size:13px;font-weight:600;color:#C0C4D8;">{len(known)} classified</div></div>
            <div><div style="font-family:'JetBrains Mono',monospace;font-size:9px;letter-spacing:.12em;color:#404560;text-transform:uppercase;margin-bottom:3px;">Size</div>
                 <div style="font-family:'Syne',sans-serif;font-size:13px;font-weight:600;color:#C0C4D8;">{len(pptx_bytes)//1024} KB</div></div>
        </div>
        """, unsafe_allow_html=True)


# Footer
st.markdown("""
<div style="margin-top:4rem;padding-top:1.5rem;border-top:1px solid rgba(255,255,255,.05);display:flex;justify-content:space-between;flex-wrap:wrap;gap:.5rem;">
    <div style="font-family:'JetBrains Mono',monospace;font-size:10px;color:#252840;">DELTAV ARCHITECTURE GENERATOR · EMERSON AUTOMATION SOLUTIONS · PUNE</div>
    <div style="font-family:'JetBrains Mono',monospace;font-size:10px;color:#252840;">RULE-BASED ENGINE · 100% OFFLINE · IMAGE EDITION v3.0</div>
</div>
""", unsafe_allow_html=True)